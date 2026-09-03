"""Generates the Account Health platform training deck (PowerPoint) from structured content.

Run: python3 training/generate_training_deck.py
Output: training/Account-Health-Training.pptx

Re-run this script any time the app changes to regenerate the deck instead of hand-editing
the binary .pptx (edit the SLIDES data below, then re-run).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand palette lifted from the app's own Tailwind theme (slate header + sky/emerald accents).
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_800 = RGBColor(0x1E, 0x29, 0x3B)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_600 = RGBColor(0x47, 0x55, 0x69)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SKY_600 = RGBColor(0x02, 0x84, 0xC7)
EMERALD_600 = RGBColor(0x05, 0x96, 0x69)
AMBER_600 = RGBColor(0xD9, 0x77, 0x06)
ROSE_600 = RGBColor(0xE1, 0x1D, 0x48)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_bg(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)
    return shape


def add_textbox(slide, left, top, width, height, text, size=18, color=SLATE_800, bold=False,
                 align=PP_ALIGN.LEFT, font="Calibri", anchor=None, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


def add_kicker_title(slide, kicker, title, kicker_color=SKY_600, title_color=SLATE_900):
    add_textbox(slide, Inches(0.6), Inches(0.28), Inches(12), Inches(0.4), kicker.upper(),
                size=13, color=kicker_color, bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.62), Inches(12.1), Inches(0.9), title,
                size=30, color=title_color, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.28), Inches(2.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = kicker_color
    line.line.fill.background()


def add_bullets(slide, left, top, width, height, items, size=16, color=SLATE_700, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.line_spacing = line_spacing
        p.space_after = Pt(8 if level == 0 else 4)
        bullet = "▸ " if level == 0 else "•  "
        indent = Inches(0) if level == 0 else Inches(0.35)
        run = p.add_run()
        run.text = f"{bullet}{text}"
        run.font.size = Pt(size if level == 0 else size - 2)
        run.font.color.rgb = color if level == 0 else SLATE_600
        run.font.bold = level == 0
        p.add_run().text = ""
        # apply simple indent via paragraph_format-like left margin trick
        box.text_frame.paragraphs[-1].alignment = PP_ALIGN.LEFT
    return box


def bullet_slide(kicker, title, items, accent=SKY_600, notes=None, footer=None):
    slide = add_slide()
    add_bg(slide, WHITE)
    add_kicker_title(slide, kicker, title, kicker_color=accent)
    add_bullets(slide, Inches(0.65), Inches(1.6), Inches(12.0), Inches(5.5), items)
    if footer:
        add_textbox(slide, Inches(0.6), Inches(7.05), Inches(12.0), Inches(0.35), footer,
                    size=11, color=SLATE_600, italic=True)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def two_col_slide(kicker, title, left_title, left_items, right_title, right_items,
                   accent=SKY_600, left_color=SKY_600, right_color=EMERALD_600, notes=None):
    slide = add_slide()
    add_bg(slide, WHITE)
    add_kicker_title(slide, kicker, title, kicker_color=accent)
    col_w = Inches(5.85)
    # left card
    card_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.55), col_w, Inches(5.55))
    card_l.fill.solid(); card_l.fill.fore_color.rgb = SLATE_100; card_l.line.color.rgb = left_color; card_l.line.width = Pt(1.25)
    add_textbox(slide, Inches(0.85), Inches(1.75), col_w - Inches(0.5), Inches(0.5), left_title, size=18, bold=True, color=left_color)
    add_bullets(slide, Inches(0.85), Inches(2.3), col_w - Inches(0.5), Inches(4.6), left_items, size=14)
    # right card
    card_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.55), col_w, Inches(5.55))
    card_r.fill.solid(); card_r.fill.fore_color.rgb = SLATE_100; card_r.line.color.rgb = right_color; card_r.line.width = Pt(1.25)
    add_textbox(slide, Inches(7.1), Inches(1.75), col_w - Inches(0.5), Inches(0.5), right_title, size=18, bold=True, color=right_color)
    add_bullets(slide, Inches(7.1), Inches(2.3), col_w - Inches(0.5), Inches(4.6), right_items, size=14)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def table_slide(kicker, title, headers, rows, accent=SKY_600, col_widths=None, notes=None, font_size=13):
    slide = add_slide()
    add_bg(slide, WHITE)
    add_kicker_title(slide, kicker, title, kicker_color=accent)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    top = Inches(1.55)
    left = Inches(0.55)
    width = Inches(12.25)
    height = Inches(5.55)
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            table.columns[i].width = Emu(int(width * (w / total)))
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = SLATE_800
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(font_size)
                r.font.color.rgb = WHITE
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else SLATE_100
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(font_size - 1)
                    run.font.color.rgb = SLATE_700
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def section_slide(number, title, subtitle):
    slide = add_slide()
    add_bg(slide, SLATE_900)
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(2.5), Inches(1.2), number, size=64, color=SKY_600, bold=True)
    add_textbox(slide, Inches(0.85), Inches(3.7), Inches(11.5), Inches(1.3), title, size=36, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.9), Inches(4.6), Inches(11), Inches(0.8), subtitle, size=16, color=SLATE_100)
    return slide


# ---------------------------------------------------------------------------
# 1. TITLE SLIDE
# ---------------------------------------------------------------------------
slide = add_slide()
add_bg(slide, SLATE_900)
stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.55), SLIDE_W, Inches(0.06))
stripe.fill.solid(); stripe.fill.fore_color.rgb = SKY_600; stripe.line.fill.background()
add_textbox(slide, Inches(0.9), Inches(2.55), Inches(11.5), Inches(0.5), "ENGINEERING TRAINING", size=16, color=SKY_600, bold=True)
add_textbox(slide, Inches(0.85), Inches(3.0), Inches(11.7), Inches(1.5), "Account Health & Growth Dashboard", size=44, color=WHITE, bold=True)
add_textbox(slide, Inches(0.9), Inches(4.75), Inches(11.5), Inches(0.6),
            "Architecture, Feature Walkthrough & Implementation Deep Dive", size=20, color=SLATE_100)
add_textbox(slide, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5),
            "React + TypeScript + Vite  |  FastAPI + Python  |  Akamai NetStorage, Google Sheets, CrUX, Grover", size=13, color=SLATE_600)

# ---------------------------------------------------------------------------
# 2. AGENDA
# ---------------------------------------------------------------------------
bullet_slide(
    "Agenda", "What We'll Cover Today",
    [
        "Product overview & who this dashboard serves",
        "Tech stack, repository layout & architecture",
        "Frontend walkthrough: Summary page, Account Detail page, and the \"Matrix family\"",
        "Real-time background jobs (SSE) — how long-running scans stream progress to the UI",
        "Deep dive: the new Archive(s) picker & NetStorage snapshot browsing",
        "NetStorage data flow — CP codes, LIVE vs archive paths, credentials",
        "Other integrations: Google Sheets/Docs, Akamai identity/AppSec, CrUX/PSI, Grover",
        "Environment configuration, Docker deployment & CI/CD status",
        "Hands-on: how to add a brand-new Matrix feature",
    ],
    accent=SKY_600,
)

# ---------------------------------------------------------------------------
# 3. PRODUCT OVERVIEW
# ---------------------------------------------------------------------------
bullet_slide(
    "Product Overview", "What Is Account Health?",
    [
        "Internal stakeholder dashboard summarizing account health across Akamai-managed web properties",
        "Started as a 2-page app (Summary + Account Detail) and grew into a full \"health diagnosis\" suite",
        "Answers questions like: Is this account's traffic healthy? Are hostnames covered by security configs? "
        "How are Core Web Vitals trending? Which properties have risky feature configs?",
        ("Summary Page — KPI strip, account list with renewal-risk/expansion signals, risk/opportunity panels", 1),
        ("Account Detail Page — hero metrics, highlights, recommended actions, and a \"tool box\" of health widgets", 1),
        ("Matrix pages — deep, filterable tables + summaries + scorecards per technical dimension "
         "(features, hostnames, traffic, performance, security)", 1),
        "Audience: internal account teams / stakeholders, not end customers — no auth layer, deployed behind internal network access",
    ],
    accent=SKY_600,
)

# ---------------------------------------------------------------------------
# 4. TECH STACK
# ---------------------------------------------------------------------------
table_slide(
    "Tech Stack", "Technology Stack At a Glance",
    ["Layer", "Technology", "Notes"],
    [
        ["Frontend", "React 19 + TypeScript + Vite", "SPA, client-side routing via react-router-dom v7"],
        ["Styling", "Tailwind CSS", "Utility classes; brand palette = slate + sky/emerald/rose/amber accents"],
        ["Frontend state", "React Context + component state", "No Redux — ArchiveContext for global archive selection"],
        ["Backend", "FastAPI (Python 3.12)", "Single backend/main.py app; all routes under /api"],
        ["Background jobs", "In-process thread pool + SSE", "job_manager.py — no external queue/broker"],
        ["Data sources", "NetStorage, Google Sheets/Docs, Akamai APIs, CrUX/PSI, Grover", "Multiple integration seams, see later slides"],
        ["Packaging", "Docker (multi-stage)", "Node build stage → Python runtime stage, single container"],
        ["Dev tooling", "ESLint + typescript-eslint, tsc -b", "npm run lint / npm run build"],
    ],
    accent=SKY_600,
    col_widths=[2, 4, 6.25],
)

# ---------------------------------------------------------------------------
# 5. ARCHITECTURE
# ---------------------------------------------------------------------------
slide = add_slide()
add_bg(slide, WHITE)
add_kicker_title(slide, "Architecture", "High-Level System Architecture")


def box(slide, x, y, w, h, text, fill, text_color=WHITE, size=13, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = fill
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = text_color
    return shp


def arrow(slide, x, y, w, h, rotation=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW if rotation == 0 else MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = SLATE_600
    shp.line.fill.background()
    return shp


box(slide, Inches(0.6), Inches(1.75), Inches(2.6), Inches(1.0), "Browser\nReact SPA (Vite build)", SLATE_800)
arrow(slide, Inches(3.35), Inches(2.05), Inches(0.55), Inches(0.4))
box(slide, Inches(4.05), Inches(1.75), Inches(2.9), Inches(1.0), "FastAPI Backend\n/api/* routes", SKY_600)
arrow(slide, Inches(7.1), Inches(2.05), Inches(0.55), Inches(0.4))
box(slide, Inches(7.8), Inches(0.9), Inches(4.9), Inches(0.6), "Akamai NetStorage — CSV/JSON reports (LIVE + archive snapshots)", EMERALD_600, size=12)
box(slide, Inches(7.8), Inches(1.6), Inches(4.9), Inches(0.6), "Google Sheets/Docs — legacy mock-replacement data seam", EMERALD_600, size=12)
box(slide, Inches(7.8), Inches(2.3), Inches(4.9), Inches(0.6), "Akamai Identity + AppSec APIs — hostname coverage (EdgeGrid auth)", EMERALD_600, size=12)
box(slide, Inches(7.8), Inches(3.0), Inches(4.9), Inches(0.6), "Google CrUX + PageSpeed Insights — Core Web Vitals", EMERALD_600, size=12)
box(slide, Inches(7.8), Inches(3.7), Inches(4.9), Inches(0.6), "Grover Security Trends API — WAF/security time series", EMERALD_600, size=12)

box(slide, Inches(0.6), Inches(3.35), Inches(2.6), Inches(1.0), "In-memory Job Manager\n(threads + SSE queues)", SLATE_700)
arrow(slide, Inches(1.9), Inches(2.9), Inches(0.4), Inches(0.4), rotation=1)
box(slide, Inches(4.05), Inches(3.35), Inches(2.9), Inches(1.0), "Static SPA hosting\n(FastAPI serves dist/)", SLATE_700)
arrow(slide, Inches(5.35), Inches(2.9), Inches(0.4), Inches(0.4), rotation=1)

add_bullets(slide, Inches(0.6), Inches(4.85), Inches(12.1), Inches(2.3), [
    "Single Docker image: Node build stage compiles the React app into /dist; Python runtime stage serves the API and the compiled SPA together",
    "Frontend talks only to the FastAPI backend — no direct browser calls to NetStorage/Akamai/Grover (keeps API keys & NS credentials server-side)",
    "Long-running scans (Matrix jobs, security charts) run in background threads and stream progress back over Server-Sent Events (SSE)",
    "Data-source resolution follows a fallback chain per page: NetStorage (preferred) → Google Sheets → bundled mock data",
], size=14)

# ---------------------------------------------------------------------------
# 6. REPO STRUCTURE
# ---------------------------------------------------------------------------
two_col_slide(
    "Codebase Tour", "Repository Structure",
    "Frontend — src/",
    [
        "pages/ — one file per route (20 pages)",
        "components/ — DashboardLayout, tables, charts, dropdowns, tone.ts (status color mapping)",
        "context/ArchiveContext.tsx — global archive/NS-context state",
        "services/ — one *Jobs.ts client per Matrix feature + netstorageData.ts, googleData.ts, sseJobClient.ts",
        "data/mockData.ts — bundled fallback dataset",
        "types/dashboard.ts — shared TS interfaces for all API payloads",
    ],
    "Backend — backend/",
    [
        "main.py — FastAPI app: every route, CORS, static SPA hosting",
        "data_service.py (~2500 lines) — all data-fetching/business logic: NetStorage, Google, Akamai, CrUX, Grover",
        "job_manager.py — Job class + JobManager (thread + SSE pub/sub)",
        "mock_data.py — fallback dataset for backend endpoints",
        "account_id_map.json — accountKey → Akamai accountId / csvAccountDir mapping",
        "authHeadersNS.py / cruxVis.py — standalone NetStorage & CrUX prototyping scripts (not imported by the app)",
    ],
)

# ---------------------------------------------------------------------------
# 7. ROUTING MAP
# ---------------------------------------------------------------------------
table_slide(
    "Frontend", "Client-Side Routing Map (App.tsx)",
    ["Route", "Page Component", "Purpose"],
    [
        ["/", "SummaryPage", "KPI strip, account list, risk/opportunity panels"],
        ["/account/:accountId", "AccountDetailPage", "Hero metrics, highlights, actions, pillars, tool box links"],
        ["/account/:id/featureMatrix(/summary|/scoreCard)", "FeatureMatrix* (3)", "Property × feature enablement matrix"],
        ["/account/:id/hostmatrix/cname(/summary)", "HostMatrixCname* (2)", "Hostname → CNAME coverage"],
        ["/account/:id/trafficMatrix(/summary|/scoreCard)", "TrafficMatrix* (3)", "Edge vs origin hits per hostname"],
        ["/account/:id/perfMatrixTopN(/summary|/scoreCard)", "PerfMatrixTopN* (3)", "Core Web Vitals for top-10 hostnames"],
        ["/account/:id/secHostCoverageMatrix(/summary|/scoreCard)", "SecHostCoverageMatrix* (3)", "WAF/security-config host coverage"],
        ["/account/:id/wsaAlertMatrix(/summary|/scoreCard)", "WsaAlertMatrix* (3)", "WSA alert config × feature matrix"],
        ["/account/:id/securityFeatureCharts", "SecurityFeatureChartsPage", "Grover security trend line charts"],
        ["*", "Navigate → /", "Fallback redirect for unknown paths"],
    ],
    accent=SKY_600,
    col_widths=[4.2, 2.6, 5.5],
    font_size=12,
)

# ---------------------------------------------------------------------------
# 8. DESIGN SYSTEM
# ---------------------------------------------------------------------------
bullet_slide(
    "Frontend Foundations", "Shared UI Building Blocks",
    [
        "DashboardLayout — wraps every page: gradient header, back-to-dashboard link, the global Archive(s) picker, page title + optional account-owner chip",
        "tone.ts — centralizes status→color mapping (healthy / watch / risk / neutral) reused by MetricTiles, AccountsTable, highlights & actions lists",
        "MetricTiles — reusable KPI tile grid (used on both Summary and Account Detail pages)",
        "HealthWidgetLink — the animated \"tool box\" cards on Account Detail linking into each Matrix feature (SVG micro-animations per feature category)",
        ("AccountSearchDropdown / ArchiveSearchDropdown — same searchable-combobox UX pattern reused for two different pickers "
         "(pick an account vs. pick an archive snapshot)", 0),
        "Every Matrix table page shares one layout: Scan Progress panel (status + progress bar + streaming log) → download-CSV link → filterable data table",
        "Column-filter dropdowns, multi-select property/hostname filters, and deep-link pre-selection (via the URL's last path segment) are consistent across all six Matrix families",
    ],
    accent=SKY_600,
)

# ---------------------------------------------------------------------------
# 9. SUMMARY PAGE
# ---------------------------------------------------------------------------
bullet_slide(
    "Page Walkthrough", "Summary Page (\"/\")",
    [
        "Entry point for stakeholders — answers \"how are all my accounts doing at a glance?\"",
        ("Loads data via fetchNsSummaryDashboardData(archive) → NetStorage all_accounts_summary.json "
         "(LIVE or a selected archive snapshot)", 0),
        ("Falls back to fetchSummaryDashboardData() (Google Sheets) if NetStorage is unavailable, "
         "then shows a \"Data source\" banner explaining which path served the page", 0),
        "MetricTiles — top-line KPI strip (e.g. total accounts, at-risk count, growth ops)",
        "AccountsTable — sortable/status-badged list of every account, links into Account Detail",
        "PanelColumns — risk/opportunity summary panels",
        "AccountSearchDropdown — type-ahead search over the live account_mapping.json to jump straight to an account",
        "Archive(s) picker (in DashboardLayout) controls which NetStorage snapshot (LIVE or archive/<date>) this page — "
        "and every page reachable from it — reads from",
    ],
    accent=SKY_600,
)

# ---------------------------------------------------------------------------
# 10. ACCOUNT DETAIL PAGE
# ---------------------------------------------------------------------------
bullet_slide(
    "Page Walkthrough", "Account Detail Page (\"/account/:accountId\")",
    [
        "Deep dive for a single account: hero metrics, health highlights, recommended actions, and 5 operational \"pillars\"",
        "Data source: fetchNsAccountDashboardData(accountId, archive) → NetStorage account_<csvAccountDir>_summary.json, "
        "same NetStorage-first/Google-fallback pattern as the Summary page",
        "Hostname coverage widget: live call to fetchAccountHostnameCoverage — hits the real Akamai AppSec API (not NetStorage) via EdgeGrid auth",
        "\"Health Diagnosis & Surgical Tool Box\" section: HealthWidgetLink cards grouped by theme "
        "(Hostname/CNAME pulse, Feature heartbeat, Traffic DNA scan, Performance stethoscope, Security pulse, WSA alert scan) "
        "linking into the six Matrix feature families",
        "Archive selection made on the Summary page (or anywhere) is preserved automatically when navigating here — "
        "no manual re-entry of the archive value",
    ],
    accent=SKY_600,
)

# ---------------------------------------------------------------------------
# 11. MATRIX FAMILY PATTERN
# ---------------------------------------------------------------------------
slide = add_slide()
add_bg(slide, WHITE)
add_kicker_title(slide, "Design Pattern", "The \"Matrix Family\" — One Pattern, Six Features")
add_bullets(slide, Inches(0.65), Inches(1.55), Inches(12.0), Inches(1.6), [
    "Every technical dimension (features, hostnames, traffic, performance, security host coverage, WSA alerts) ships as the SAME 3-page shape:",
])
box(slide, Inches(0.8), Inches(3.2), Inches(3.6), Inches(1.1), "1) Table Page\nFull filterable row-level data + CSV download", SKY_600, size=13)
arrow(slide, Inches(4.55), Inches(3.55), Inches(0.6), Inches(0.4))
box(slide, Inches(5.3), Inches(3.2), Inches(3.6), Inches(1.1), "2) Summary Page\nAggregated totals + breakdown pie/line charts", EMERALD_600, size=13)
arrow(slide, Inches(9.05), Inches(3.55), Inches(0.6), Inches(0.4))
box(slide, Inches(9.8), Inches(3.2), Inches(2.7), Inches(1.1), "3) ScoreCard Page\nPer-entity grouped JSON view + raw JSON toggle", AMBER_600, size=13)
add_bullets(slide, Inches(0.65), Inches(4.7), Inches(12.0), Inches(2.4), [
    "All three pages for a feature share one background job pipeline (CSV → parse → transform), started once and re-used via cross-links",
    "All three read the SAME data mode (csv_data_remote — always NetStorage) and the SAME inherited Context (NS base path)",
    "Consistency payoff: once you understand ONE Matrix family (e.g. Feature Matrix), you understand all six — "
    "the only thing that changes is the CSV schema and the domain-specific columns",
], size=15)

# ---------------------------------------------------------------------------
# 12. MATRIX CATALOG
# ---------------------------------------------------------------------------
table_slide(
    "Feature Catalog", "The Six Matrix Families",
    ["Feature", "Source CSV", "Answers…"],
    [
        ["Feature Matrix", "config-audit.csv", "Which Akamai properties have which features enabled?"],
        ["Hostname CNAME Matrix", "config-summary.csv", "Are hostnames CNAME'd correctly to Akamai edge?"],
        ["Traffic Matrix", "traffic-report-hits-by-hostname.csv", "Edge vs. origin hit distribution per hostname"],
        ["Perf Matrix (Top 10)", "traffic-report-hits-by-hostname.csv + live CrUX/PSI", "Core Web Vitals for the busiest hostnames"],
        ["Sec Host Coverage Matrix", "hostname-coverage-*.csv", "Are hostnames covered by a WAF/security config?"],
        ["WSA Alert Matrix", "wsa-alert-*.csv", "Which configs have which WSA alert features enabled?"],
    ],
    accent=EMERALD_600,
    col_widths=[3, 4.5, 6.5],
)

# ---------------------------------------------------------------------------
# 13. FEATURE MATRIX DEEP DIVE
# ---------------------------------------------------------------------------
bullet_slide(
    "Deep Dive", "Feature Matrix — Implementation Notes",
    [
        "Backend: get_account_feature_matrix() downloads config-audit.csv from NetStorage, parses rows into property × feature columns",
        "Frontend: runFeatureMatrixJob() streams progress via SSE while the CSV downloads/parses server-side",
        "Multi-select property filter + per-column value filters (checkbox dropdown) combine client-side over the already-fetched rows",
        "Deep-link support: /featureMatrix/:propIdOrFeature pre-selects a property OR pre-filters a feature column to \"Enabled\" "
        "— used by HealthWidgetLink and cross-links from other pages",
        "Summary variant aggregates Enabled/Disabled totals + per-column breakdowns for a pie chart",
        "ScoreCard variant regroups by feature name → property count, with an expandable raw-JSON view for power users/debugging",
        "\"Download CSV\" link re-downloads the exact same file NetStorage served (no caching — always fresh)",
    ],
    accent=ROSE_600,
)

# ---------------------------------------------------------------------------
# 14. PERF MATRIX TOPN DEEP DIVE
# ---------------------------------------------------------------------------
bullet_slide(
    "Deep Dive", "Performance Matrix (Top 10) — CrUX + PageSpeed Insights",
    [
        "Cost-control design: live CrUX/PSI lookups are expensive, so this feature only tests the top 10 hostnames "
        "by 7-day edge hits (from the Traffic Matrix CSV) instead of every hostname",
        "get_hostname_core_web_vitals(): tries Chrome UX Report (CrUX) records:queryRecord first (real user field data)",
        ("On CrUX 404 (insufficient real-user traffic) → falls back to PageSpeed Insights runPagespeed "
         "(synthetic Lighthouse lab run, mobile strategy)", 0),
        "Metrics: LCP, INP, CLS — each classified good / needs-improvement / poor against Google's published thresholds",
        "Gotcha we hit in production: CrUX returns LCP/INP percentiles as numbers but CLS as a numeric STRING — "
        "_coerce_metric_value() safely float()-casts every metric instead of assuming type",
        "Hostnames are fetched concurrently via a ThreadPoolExecutor (default 5 workers) with live progress logged per hostname",
        "History endpoint (queryHistoryRecord) is fetched best-effort for the trend line chart; failures there don't fail the whole request",
    ],
    accent=ROSE_600,
)

# ---------------------------------------------------------------------------
# 15. SECURITY FEATURE CHARTS
# ---------------------------------------------------------------------------
bullet_slide(
    "Deep Dive", "Security Feature Charts — Grover Integration",
    [
        "The only page NOT following the Matrix 3-page pattern — a single interactive chart page instead",
        "User picks a start/end date range + account name, then triggers a background job",
        "Backend calls Grover's security-trends API (api.grover.akamai.com) with an X-API-Key header (X_API_KEY env var)",
        "extract_security_trend_series() normalizes whatever shape the API returns (bare list, or wrapped in data/results/records/trends/series) "
        "into per-dimension time series, auto-detecting the date field from a candidate-key list",
        "Same defensive numeric coercion pattern as CrUX: metric values may arrive as numeric strings",
        "Frontend renders each numeric dimension as a toggle-able line on SecurityTrendLineChart; unrecognized shapes still show the raw JSON so nothing is silently lost",
    ],
    accent=ROSE_600,
)

# ---------------------------------------------------------------------------
# 16. BACKGROUND JOBS / SSE ARCHITECTURE
# ---------------------------------------------------------------------------
slide = add_slide()
add_bg(slide, WHITE)
add_kicker_title(slide, "Real-Time Updates", "Background Jobs & Server-Sent Events (SSE)")
box(slide, Inches(0.6), Inches(1.7), Inches(2.6), Inches(0.9), "1. POST .../jobs\nCreate + start job", SKY_600, size=12)
arrow(slide, Inches(3.3), Inches(1.95), Inches(0.5), Inches(0.4))
box(slide, Inches(3.95), Inches(1.7), Inches(2.7), Inches(0.9), "2. job_manager\nruns target() on a\ndaemon thread", SLATE_700, size=12)
arrow(slide, Inches(6.75), Inches(1.95), Inches(0.5), Inches(0.4))
box(slide, Inches(7.4), Inches(1.7), Inches(2.7), Inches(0.9), "3. Job.log()/complete()\npush to subscriber\nqueues", EMERALD_600, size=12)
arrow(slide, Inches(6.6), Inches(2.7), Inches(0.4), Inches(0.4), rotation=1)
box(slide, Inches(3.95), Inches(3.2), Inches(4.2), Inches(0.9), "4. GET .../jobs/{id}/events\ntext/event-stream (SSE)", AMBER_600, size=12)
arrow(slide, Inches(3.7), Inches(3.4), Inches(0.4), Inches(0.4))
box(slide, Inches(0.6), Inches(3.2), Inches(2.9), Inches(0.9), "5. runJobWithRetry()\nparses SSE, updates\nReact state live", SKY_600, size=12)
add_bullets(slide, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.6), [
    "Job (job_manager.py): tracks status (running/completed/failed), percent, message history, and a list of subscriber Queues",
    "Every .log() call appends to history AND pushes to any live subscribers — new SSE connections get instant replay of everything so far",
    "Client-side runJobWithRetry() (sseJobClient.ts) is the single retry engine used by ALL six Matrix families: "
    "auto-retries the whole start+stream sequence on 404/connection failure with exponential backoff (capped at 15s)",
    "Why 404 matters: jobs live in memory only — a backend restart mid-scan loses in-flight jobs, so the client transparently restarts the job rather than getting stuck",
    "Every Matrix feature also exposes a synchronous, non-SSE JSON variant of summary/scoreCard endpoints for simple direct consumption",
], size=14)

# ---------------------------------------------------------------------------
# 17. SECTION DIVIDER — Archive Feature
# ---------------------------------------------------------------------------
section_slide("★", "Spotlight Feature: Archive(s) & NS Context Picker",
              "A new global, persisted control for browsing historical NetStorage snapshots across every page")

# ---------------------------------------------------------------------------
# 18. ARCHIVE FEATURE OVERVIEW
# ---------------------------------------------------------------------------
bullet_slide(
    "New Feature", "Archive(s) Picker — What It Does",
    [
        "Problem it replaced: every page had a free-text \"Archive(s):\" box, plus every Matrix page had ITS OWN separately-typed "
        "\"Context (NS base path):\" box — easy to get out of sync, and the value never carried over between pages",
        "New behavior: one searchable dropdown/combobox in the header (DashboardLayout), fed by the REAL list of archive folders in NetStorage",
        "Selecting an archive is now a single global action that persists across the entire app — Summary, Account Detail, "
        "and all 18 Matrix/Summary/ScoreCard pages automatically inherit it",
        "Default value is LIVE (staticSiteContent) until the user explicitly picks a snapshot",
        "\"Context (NS base path):\" on every Matrix page is now READ-ONLY — it always mirrors the global Archive(s) selection; "
        "users only ever change it in one place",
        "CSV/JSON fetch logic itself was intentionally left untouched — this feature only changes WHERE the context value comes from, "
        "not how it's used to fetch data",
    ],
    accent=AMBER_600,
)

# ---------------------------------------------------------------------------
# 19. ARCHIVE FEATURE DATA FLOW
# ---------------------------------------------------------------------------
slide = add_slide()
add_bg(slide, WHITE)
add_kicker_title(slide, "New Feature", "Archive(s) Picker — Data Flow", kicker_color=AMBER_600)
box(slide, Inches(0.6), Inches(1.7), Inches(3.0), Inches(1.0), "GET /api/dashboard/ns/archives\n(backend, cached 60s)", SKY_600, size=12)
arrow(slide, Inches(3.75), Inches(2.0), Inches(0.5), Inches(0.4))
box(slide, Inches(4.4), Inches(1.7), Inches(3.1), Inches(1.0), "ArchiveContext\n(React context, app-wide)", SLATE_700, size=12)
arrow(slide, Inches(7.65), Inches(2.0), Inches(0.5), Inches(0.4))
box(slide, Inches(8.3), Inches(1.7), Inches(3.6), Inches(1.0), "ArchiveSearchDropdown\nin DashboardLayout header", EMERALD_600, size=12)
arrow(slide, Inches(2.0), Inches(2.75), Inches(0.4), Inches(0.4), rotation=1)
box(slide, Inches(0.6), Inches(3.3), Inches(3.0), Inches(1.0), "localStorage +\ncurrent route's ?archive=", AMBER_600, size=12)
arrow(slide, Inches(9.9), Inches(2.75), Inches(0.4), Inches(0.4), rotation=1)
box(slide, Inches(8.3), Inches(3.3), Inches(3.6), Inches(1.0), "Every page: useArchive()\n→ contextPath (read-only display)", AMBER_600, size=12)
add_bullets(slide, Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.5), [
    "ArchiveProvider wraps <Routes> (inside BrowserRouter) → survives client-side navigation without re-fetching or resetting",
    "State also mirrors to localStorage (survives full page reloads, e.g. the account-search full-navigation flow) "
    "and to the current route's ?archive= query param (shareable URLs)",
    "Every Matrix/Summary/ScoreCard page just calls useArchive() and reads contextPath instead of maintaining its own local state — "
    "17 pages refactored to this one pattern",
    "Backend: list_ns_archive_folders() lists CPCODE/archive via the NetStorage \"dir\" action, returns [\"archive/20260901\", ...] newest-first",
], size=14)

# ---------------------------------------------------------------------------
# 20. ARCHIVE FEATURE — THE BUG WE FIXED
# ---------------------------------------------------------------------------
bullet_slide(
    "New Feature — Lesson Learned", "Production Bug: NetStorage \"Implicit Directories\"",
    [
        "Symptom: GET /api/dashboard/ns/archives returned an empty list with "
        "\"NetStorage stat failed for '/CPCODE/archive': status=404, reason='Not Found'\" — even though archive/20260810, "
        "/20260819, /20260901 clearly existed and were being used successfully for CSV/JSON fetches",
        "Root cause: NetStorage's stat action only returns information about the exact object requested. "
        "Directories like /CPCODE/archive are \"implicit\" — they have no real directory-marker object, they're inferred purely "
        "from the paths of files underneath them — so stat() 404s on them",
        "The dir (or list) action, by contrast, lists a path's children by prefix and works fine on implicit directories",
        "Fix: switched netstorage.stat(remote_dir) → netstorage.dir(remote_dir) in list_ns_archive_folders()",
        "Verified live against the real NetStorage account: confirmed real layout is /CPCODE/archive/<YYYYMMDD>/... "
        "alongside /CPCODE/staticSiteContent/... (LIVE) and /CPCODE/demo/...",
        "Takeaway for the team: when scripting against NetStorage, use dir/list to enumerate a folder's contents — "
        "reserve stat for confirming a single known object exists",
    ],
    accent=AMBER_600,
)

# ---------------------------------------------------------------------------
# 21. NETSTORAGE OVERVIEW
# ---------------------------------------------------------------------------
section_slide("2", "NetStorage Data Flow", "How every report CSV/JSON actually gets from Akamai storage to the browser")

# ---------------------------------------------------------------------------
# 22. NETSTORAGE PATH MODEL
# ---------------------------------------------------------------------------
bullet_slide(
    "NetStorage Deep Dive", "The Path Model: CP Code + Base Path",
    [
        "Every NetStorage object lives at /{NS_CP_CODE}/{base_path}/{relative_path}",
        "NS_CP_CODE — the Akamai NetStorage CP code (account-specific numeric ID), e.g. 2052217",
        "base_path — normally NS_BASE_PATH env var (\"staticSiteContent\" = the LIVE folder); can be OVERRIDDEN per-request by a \"context\" query param",
        "Real top-level layout for this account (confirmed live): staticSiteContent/ (LIVE), archive/ (dated snapshots), demo/",
        "context=\"archive/20260901\" → base_path becomes \"archive/20260901\" for that one request only — nothing else in the app changes",
        "context=None/empty (\"\" ) → falls back to NS_BASE_PATH (LIVE) — this is the default on every page until a user picks an archive",
        "This single override mechanism is what powers BOTH the old free-text box AND the new Archive(s) dropdown — "
        "the UI changed, the underlying fetch contract did not",
    ],
    accent=EMERALD_600,
)

# ---------------------------------------------------------------------------
# 23. NETSTORAGE FETCH PIPELINE
# ---------------------------------------------------------------------------
bullet_slide(
    "NetStorage Deep Dive", "CSV/JSON Fetch Pipeline",
    [
        "download_csv_from_netstorage(remote_path, local_path) — the one low-level function every fetch goes through: "
        "builds an akamai.netstorage.Netstorage client from env creds, calls .download(), verifies the file actually landed on disk",
        "Rich diagnostics on failure: captures status code, reason, response headers/body snippet, and adds a hint for common codes "
        "(404 → wrong path/CP code, 403 → bad credentials/ACL)",
        "resolve_report_csv_path() — per-account CSV reports (config-audit.csv, traffic CSVs, etc.), keyed by csvAccountDir from account_id_map",
        "download_ns_json() — JSON snapshots (all_accounts_summary.json, account_<dir>_summary.json, account_mapping.json)",
        "Design decision: NO local caching for report downloads — every request re-downloads fresh from NetStorage, "
        "so the dashboard never shows stale data (account-mapping and archive-list ARE cached briefly: 30s / 60s TTL, since they change rarely)",
        "Every one of the 6 Matrix families' job functions calls resolve_report_csv_path() with its own relative CSV path constant",
    ],
    accent=EMERALD_600,
)

# ---------------------------------------------------------------------------
# 24. NETSTORAGE SDK / CREDS
# ---------------------------------------------------------------------------
table_slide(
    "NetStorage Deep Dive", "SDK Actions & Required Credentials",
    ["NS SDK Action", "Use in this app", "Works on implicit dirs?"],
    [
        ["download(remote, local)", "Every CSV/JSON fetch (download_csv_from_netstorage)", "N/A"],
        ["dir(path)", "list_ns_archive_folders() — list archive/<date> folders", "✅ Yes"],
        ["list(path)", "Not currently used (recursive listing) — available if needed", "✅ Yes"],
        ["stat(path)", "Not used for directory listing (see incident, prior slide)", "❌ No — 404s"],
        ["upload / mkdir / delete", "Not used by this app (read-only integration)", "—"],
    ],
    accent=EMERALD_600,
    col_widths=[3.5, 6, 3],
)
bullet_slide(
    "NetStorage Deep Dive", "Required Environment Variables",
    [
        "NS_HOSTNAME — NetStorage upload domain (e.g. mysitelab-nsu.akamaihd.net)",
        "NS_KEYNAME / NS_KEY — NetStorage API key credentials (HMAC-signed requests, never sent to the browser)",
        "NS_CP_CODE — the CP code that scopes every remote path",
        "NS_BASE_PATH — the LIVE base path segment (staticSiteContent)",
        "All read via get_ns_config() — a single function, so there's one place to see every NetStorage-related setting",
        "Set only in .env.server (backend, server-side) — never in the frontend's .env, and never committed to git",
    ],
    accent=EMERALD_600,
)

# ---------------------------------------------------------------------------
# 26. FALLBACK CHAINS
# ---------------------------------------------------------------------------
bullet_slide(
    "Resilience Pattern", "Data-Source Fallback Chains",
    [
        "Summary & Account Detail pages: NetStorage (preferred, live/archive) → Google Sheets (fetchSummaryDashboardData / "
        "fetchAccountDashboardData) → bundled mock data, with a visible \"Data source\" banner explaining which path served the page",
        "Account mapping (load_account_id_map): NetStorage account_mapping.json → local backend/account_id_map.json fallback",
        "Perf Matrix: CrUX field data → PageSpeed Insights synthetic lab fallback on 404 (per-hostname, not page-level)",
        "Every NetStorage call wraps failures into a structured {source, data, error} response so the frontend can render "
        "a clear message instead of a blank page",
        "This layered design means the demo/dev experience (no NetStorage creds) still works end-to-end against mock data — "
        "useful for onboarding new engineers before they get production credentials",
    ],
    accent=SLATE_700,
)

# ---------------------------------------------------------------------------
# 27. GOOGLE SHEETS/DOCS
# ---------------------------------------------------------------------------
table_slide(
    "Other Integrations", "Google Sheets/Docs Integration",
    ["Sheet Tab", "Columns", "Feeds"],
    [
        ["SummaryMetrics", "id,title,value,subtitle,tone", "Summary KPI tiles"],
        ["Accounts", "accountId,name,healthScore,healthTone,renewalRisk,...", "Summary accounts table"],
        ["SummaryPanels", "panelId,panelTitle,itemId,label,value,tone", "Summary risk/opportunity panels"],
        ["AccountDetails", "accountId,name,owner,quarter", "Account Detail header"],
        ["AccountHeroMetrics / Highlights / Actions / Pillars", "see README schema", "Account Detail sections"],
    ],
    accent=SLATE_700,
    col_widths=[3.5, 6, 3],
)
bullet_slide(
    "Other Integrations", "Google Sheets/Docs — Notes",
    [
        "Backend integration: backend/data_service.py using a Google service-account credential (GOOGLE_SERVICE_ACCOUNT_EMAIL / "
        "GOOGLE_PRIVATE_KEY / GOOGLE_SHEETS_SPREADSHEET_ID)",
        "Browser-side alternative also exists (src/services/googleData.ts + VITE_GOOGLE_API_KEY) — NOT recommended for internal secrets, "
        "kept mainly for quick local prototyping",
        "docId column (optional): when present, the linked Google Doc's text OVERRIDES the sheet cell value — "
        "lets non-engineers edit longer narrative text without touching the sheet",
        "tone values are constrained to healthy / watch / risk / neutral and drive every colored badge in the UI via tone.ts",
        "This was the ORIGINAL data source before NetStorage was added — now serves purely as the fallback layer",
    ],
    accent=SLATE_700,
)

# ---------------------------------------------------------------------------
# 28. AKAMAI HOSTNAME COVERAGE
# ---------------------------------------------------------------------------
bullet_slide(
    "Other Integrations", "Akamai Hostname Coverage (Live API, not NetStorage)",
    [
        "The ONE feature in the app that calls live Akamai management APIs directly, using EdgeGrid authentication (~/.edgerc)",
        "create_akamai_session(): builds a requests.Session with EdgeGridAuth.from_edgerc(), scoped to EDGE_RC_SECTION",
        "Step 1 — Identity API: GET identity-management/v3/api-clients/self/account-switch-keys?search={accountId} "
        "resolves an accountSwitchKey for the target account",
        "resolve_account_switch_key(): handles 0/1/many matches, disambiguating on exact accountName match when needed",
        "Step 2 — AppSec API: GET appsec/v1/hostname-coverage (with accountSwitchKey applied to every subsequent request) "
        "returns covered/not_covered/unknown per hostname, plus security config + policy names",
        "Both calls are rate-limited (100 calls/60s) via the ratelimit package's @sleep_and_retry/@limits decorators",
        "Powers the Covered/Not Covered/Unknown counters on the Account Detail page's tool box header",
    ],
    accent=SLATE_700,
)

# ---------------------------------------------------------------------------
# 29. ENV CONFIG
# ---------------------------------------------------------------------------
table_slide(
    "Configuration", "Environment Variables Cheat Sheet",
    ["Variable", "File", "Purpose"],
    [
        ["VITE_APP_BASE_PATH", ".env", "Frontend route prefix for cloud hosting (e.g. /account-health/)"],
        ["VITE_DASHBOARD_DATA_MODE / VITE_API_BASE_URL", ".env", "backend vs google mode; API base URL"],
        ["APP_BASE_PATH / FRONTEND_ORIGIN", ".env.server", "Backend route prefix + CORS allow-origin"],
        ["NS_HOSTNAME / NS_KEYNAME / NS_KEY / NS_CP_CODE / NS_BASE_PATH", ".env.server", "NetStorage credentials + LIVE base path"],
        ["GOOGLE_SHEETS_SPREADSHEET_ID / GOOGLE_SERVICE_ACCOUNT_EMAIL / GOOGLE_PRIVATE_KEY", ".env.server", "Google Sheets/Docs fallback"],
        ["CRUX_API_KEY", ".env.server", "Chrome UX Report + PageSpeed Insights (perf matrix)"],
        ["X_API_KEY", ".env.server", "Grover security-trends API"],
        ["EDGE_RC_SECTION / EDGE_RC_PATH / AKAMAI_ACCOUNT_MAP_PATH", ".env.server", "Akamai EdgeGrid auth for hostname coverage"],
    ],
    accent=SLATE_700,
    col_widths=[5, 2, 6],
    font_size=12,
)

# ---------------------------------------------------------------------------
# 30. BACKEND API CONVENTIONS
# ---------------------------------------------------------------------------
bullet_slide(
    "Backend Design", "API Design Conventions (main.py)",
    [
        "One FastAPI app, one router file (main.py) — all business logic lives in data_service.py, main.py just wires HTTP → functions",
        "APP_BASE_PATH env var is normalized once (normalize_base_path) and prefixed onto every route via API_PREFIX — "
        "supports both root-hosted (\"/\") and prefixed cloud hosting (\"/account-health\") from the same code",
        "CORS restricted to a single FRONTEND_ORIGIN (default http://localhost:5173) — not a wildcard",
        "Per-Matrix-feature endpoint quartet: POST .../jobs (start, SSE-driven) · GET .../jobs/{id}/events (SSE stream) · "
        "GET .../summary or /scoreCard (synchronous JSON, no job) · GET .../download (raw CSV as file attachment)",
        "context query param (archive override) and data query param (data source mode, currently always csv_data_remote) "
        "are consistent across every Matrix endpoint",
        "Static SPA hosting: FastAPI mounts dist/assets and serves index.html for any unmatched path under APP_BASE_PATH — "
        "one container serves both the API and the built React app",
    ],
    accent=SKY_600,
)

# ---------------------------------------------------------------------------
# 31. LOCAL DEV WORKFLOW
# ---------------------------------------------------------------------------
table_slide(
    "Developer Workflow", "Local Development Commands",
    ["Command", "What it does"],
    [
        ["npm install", "Install frontend dependencies"],
        ["npm run dev", "Vite dev server only (frontend, hot reload)"],
        ["npm run server:dev", "FastAPI backend only, --reload, port 4000"],
        ["npm run dev:full", "Runs both concurrently (recommended for full-stack work)"],
        ["python3 -m venv .venv && .venv/bin/python -m pip install -r requirements.txt", "One-time backend Python env setup"],
        ["npm run build", "tsc -b (type-check) then vite build → dist/"],
        ["npm run lint", "ESLint (typescript-eslint + react-hooks + react-refresh rules)"],
        ["npm run preview", "Serve the production build locally"],
    ],
    accent=SKY_600,
    col_widths=[6, 7],
)

# ---------------------------------------------------------------------------
# 32. DOCKER
# ---------------------------------------------------------------------------
bullet_slide(
    "Deployment", "Docker: Multi-Stage Build",
    [
        "Stage 1 (node:22-alpine) — npm ci, then npm run build with a build-arg VITE_APP_BASE_PATH baked into the compiled JS "
        "(frontend env vars are compile-time, not runtime)",
        "Stage 2 (python:3.12-slim) — pip install -r requirements.txt, copies backend/ source + the built dist/ from stage 1",
        "Final image only contains: Python runtime + backend code + compiled static frontend — no Node/npm in the shipped image",
        "Runtime configuration is 100% environment variables (APP_BASE_PATH, SERVER_DATA_MODE, NS_*, GOOGLE_*, CRUX_API_KEY, X_API_KEY) — "
        "no secrets baked into the image, so the same image can be promoted across environments",
        "Container command: python -m uvicorn backend.main:app --host 0.0.0.0 --port ${PORT:-4000}",
        "Typical run: docker run --rm -p 4000:4000 --env-file .env.server -e APP_BASE_PATH=/account-health account-health:latest",
    ],
    accent=SKY_600,
)

# ---------------------------------------------------------------------------
# 33. CI/CD STATUS
# ---------------------------------------------------------------------------
bullet_slide(
    "CI/CD", "Current State & Recommendations",
    [
        "As of today, the repo has NO GitHub Actions workflows configured (.github/ only contains an editor "
        "copilot-instructions.md, no .github/workflows/*.yml)",
        "Everything currently described as \"the workflow\" is a manual/local developer workflow: npm run build, "
        "docker build, then a manual docker run / deploy to the target platform",
        "Recommended pipeline stages to add (suggested for discussion, not yet implemented):",
        ("Lint + type-check: npm run lint, tsc -b", 1),
        ("Frontend build verification: npm run build", 1),
        ("Backend smoke import / basic pytest (none exist yet — no automated backend tests today)", 1),
        ("docker build (multi-stage) as a build-verification gate on PRs", 1),
        ("Optional: push to a registry + deploy step for main-branch merges", 1),
        "Flagging this gap explicitly so new team members don't assume CI gates exist today — code review is currently the only gate",
    ],
    accent=ROSE_600,
)

# ---------------------------------------------------------------------------
# 34. SECURITY & RESILIENCE RECAP
# ---------------------------------------------------------------------------
bullet_slide(
    "Cross-Cutting Concerns", "Security & Resilience Practices Used Throughout",
    [
        "Secrets only via environment variables, loaded with python-dotenv from .env.server — never hardcoded, never committed "
        "(.env.server.example / account_id_map.example.json are the checked-in templates)",
        "NetStorage/Akamai/Grover credentials never reach the browser — every external call is proxied through the FastAPI backend",
        "CORS locked to a single configured FRONTEND_ORIGIN, not a wildcard",
        "Outbound API calls (Akamai identity/AppSec, CrUX, PageSpeed, Grover) are all rate-limited via @sleep_and_retry/@limits",
        "Defensive numeric parsing everywhere third-party JSON is involved (CrUX CLS-as-string, Grover metrics-as-string) — "
        "never a bare cast, always try/except float()",
        "Multi-layer fallback chains (NetStorage → Google → mock) mean a single integration outage degrades gracefully "
        "instead of breaking the page",
        "SSE job client auto-recovers from backend restarts (404-triggered retry with backoff) instead of leaving the UI stuck",
    ],
    accent=SLATE_700,
)

# ---------------------------------------------------------------------------
# 35. HOW TO ADD A NEW MATRIX FEATURE
# ---------------------------------------------------------------------------
bullet_slide(
    "Hands-On Guide", "How to Add a New Matrix Feature (Recipe)",
    [
        "1. Backend (data_service.py): define RELATIVE_PATH constant for the new CSV, then 3 functions — "
        "get_account_<x>_matrix / _summary / _scorecard, each calling resolve_report_csv_path() + parsing rows",
        "2. Backend (main.py): add the endpoint quartet — POST jobs, GET jobs/{id}/events, GET summary/scoreCard (sync JSON), GET download — "
        "for the table page AND the summary/scoreCard variants",
        "3. Frontend (services/<x>Jobs.ts): mirror an existing Jobs.ts file — startJob/runJob (SSE) + fetch<X>Summary/ScoreCard (sync) + download URL helper",
        "4. Frontend (types/dashboard.ts): add the result/row/progress-event interfaces for the new feature",
        "5. Frontend (pages/): copy an existing 3-page family (Page/Summary/ScoreCard) as a template — swap the service import, "
        "result type, and domain-specific columns/filters; keep the shared useArchive()/DashboardLayout/Scan-Progress pattern intact",
        "6. src/App.tsx: add the 3 new routes (table / summary / scoreCard)",
        "7. AccountDetailPage.tsx: add a new HealthWidgetLink group (or extend an existing one) so users can discover the feature",
        "Because every family follows the identical pattern, a new Matrix feature is mostly copy-adapt, not net-new architecture",
    ],
    accent=SKY_600,
)

# ---------------------------------------------------------------------------
# 36. KEY FILES CHEAT SHEET
# ---------------------------------------------------------------------------
table_slide(
    "Quick Reference", "Key Files Cheat Sheet",
    ["File", "What lives here"],
    [
        ["backend/main.py", "Every HTTP route; CORS; static SPA hosting; APP_BASE_PATH prefixing"],
        ["backend/data_service.py", "All business logic: NetStorage, Google, Akamai, CrUX/PSI, Grover, CSV parsing"],
        ["backend/job_manager.py", "Job/JobManager — background thread execution + SSE pub/sub"],
        ["src/App.tsx", "Client-side route table + ArchiveProvider mount point"],
        ["src/context/ArchiveContext.tsx", "Global archive/NS-context state (localStorage + URL sync)"],
        ["src/components/DashboardLayout.tsx", "Shared page chrome + the Archive(s) picker"],
        ["src/services/sseJobClient.ts", "Shared SSE start+stream+retry engine used by every Matrix feature"],
        ["src/services/netstorageData.ts", "Frontend client for all /api/dashboard/ns/* endpoints"],
        ["src/types/dashboard.ts", "Shared TypeScript interfaces for every API payload"],
        [".env.server / .env.server.example", "All backend runtime secrets & config (never committed with real values)"],
        ["Dockerfile", "Multi-stage build: Node build → Python runtime, single deployable image"],
    ],
    accent=SKY_600,
    col_widths=[4.5, 8.5],
    font_size=12,
)

# ---------------------------------------------------------------------------
# 37. RECAP / Q&A
# ---------------------------------------------------------------------------
slide = add_slide()
add_bg(slide, SLATE_900)
add_textbox(slide, Inches(0.85), Inches(1.2), Inches(11.5), Inches(1.0), "Recap & Questions", size=38, color=WHITE, bold=True)
add_bullets(slide, Inches(0.9), Inches(2.3), Inches(11.3), Inches(4.2), [
    "One consistent Matrix pattern (Table / Summary / ScoreCard) powers six different technical dimensions",
    "One SSE job engine + one archive-context system are shared across the entire app — learn them once",
    "NetStorage is the primary data source; Google Sheets and mock data are graceful fallbacks, never silent failures",
    "Multiple live external integrations (Akamai identity/AppSec, CrUX/PSI, Grover) all follow the same defensive patterns: "
    "rate limiting, safe numeric coercion, and clear error surfacing",
    "No CI/CD pipeline exists yet — an open opportunity for the team",
    "Adding a new Matrix feature is a well-trodden recipe, not a fresh design exercise",
], size=17, color=SLATE_100)
add_textbox(slide, Inches(0.9), Inches(6.7), Inches(11.3), Inches(0.5),
            "Resources: README.md · backend/data_service.py · src/context/ArchiveContext.tsx · this deck's generator script (training/generate_training_deck.py)",
            size=12, color=SLATE_600)

# ---------------------------------------------------------------------------
OUTPUT_PATH = "/Users/gramesh/Documents/AHDev-groverIntegration/account-health/training/Account-Health-Training.pptx"
prs.save(OUTPUT_PATH)
print(f"Saved {len(prs.slides)} slides to {OUTPUT_PATH}")
